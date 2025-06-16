import os
from openai import OpenAI
from dotenv import load_dotenv
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx import Presentation
import requests
from io import BytesIO
import time
import json
from PIL import Image

load_dotenv()

os.makedirs('powerpoint-ppt', exist_ok=True)

client = OpenAI(
    api_key=os.getenv("OPENAI_API_KEY")
)

TITLE_FONT_SIZE = Pt(30)
CONTENT_FONT_SIZE = Pt(16)

create_slide_contents_desc = [
    {
                "name": "create_slide_contents",
                "description": "Create contents for the slide based on the slide title passed as a parameter",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "slide_title": {
                            "type": "string",
                            "description": "slide title, e.g. United Nations & It's Role",
                        },
                    },
                    "required": ["slide_title"],
                },
            }
]

create_slide_title_desc = [
    {
                "name": "create_slide_title",
                "description": "Create title for the slide based on the topic & number of slides required passed as a parameter",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "topic": {
                            "type": "string",
                            "description": "slide title, e.g. United Nations & It's Role",
                        },
                        "num_slides":{
                            "type":"string",
                            "description":"Number Of Slides Required"
                        }
                    },
                    "required": ["topic","num_slides"],
                },
            }
]

create_ppt_desc = [
    {
                "name": "create_presentations",
                "description": "Create presentation for the slide based on the topic & number of slides required passed as a parameter",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "topic": {
                            "type": "string",
                            "description": "slide title, e.g. United Nations & It's Role",
                        },
                        "num_slides":{
                            "type":"string",
                            "description":"Number Of Slides Required"
                        }
                    },
                    "required": ["topic","num_slides"],
                },
            }
]


def generate_dalle_image(prompt):
    response = client.images.generate(
        model="dall-e-2",
        size="256x256",
        prompt=prompt,
        n=1,
    )
    image_url = response.data[0].url
    image_response = requests.get(image_url)
    image = Image.open(BytesIO(image_response.content))
    return image


def create_slide_title(topic, num_slides):
    prompt = f'Generate {num_slides} short slides titles for the topic {topic}'
    completion = client.chat.completions.create(
        model='gpt-3.5-turbo',
        messages=[{"role": "system", "content": prompt}],
        temperature=0.1,
        max_tokens=200
    )
    generated_content = completion.choices[0].message.content
    # Split the content into individual titles
    slide_titles = [title.strip() for title in generated_content.split('.') if title.strip()]

    return slide_titles


def create_slide_contents(slide_title):
    prompt = f'Generate content for the slide {slide_title}. The content must be in a medium-worded paragraph. Return only 1 paragraph. The paragraph should not contain more than 20 words'
    completion = client.chat.completions.create(
        model='gpt-3.5-turbo',
        messages=[{"role": "system", "content": prompt}],
        temperature=0.1,
        max_tokens=200
    )
    return completion.choices[0].message.content


def create_title_slide(powerpoint, topic):
    title_slide_layout = powerpoint.slide_layouts[0]
    title_slide = powerpoint.slides.add_slide(title_slide_layout)

    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = topic
    subtitle.text = "AI Version"

    title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title.text_frame.paragraphs[0].font.bold = True

    subtitle.text_frame.paragraphs[0].font.size = CONTENT_FONT_SIZE

    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(149, 104, 227)


def create_presentations(topic, num_slides):
    powerpoint = Presentation()

    create_title_slide(powerpoint, topic)

    title_and_content_layout = powerpoint.slide_layouts[5]

    background_color = RGBColor(149, 104, 227)

    slide_titles = create_slide_title(topic, num_slides)

    slide_titles.pop(0)

    for slide_title in slide_titles:
        slide = powerpoint.slides.add_slide(title_and_content_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        title = slide.shapes.title
        title.text = slide_title
        title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        title.text_frame.paragraphs[0].font.bold = True

        # Add text on the left side
        left_inch_text = Inches(1)  # Adjust the left margin as needed
        top_inch_text = Inches(2.3)  # Adjust the top margin as needed
        width_text = Inches(4)  # Adjust the width of the text box
        height_text = Inches(5)  # Adjust the height of the text box
        text_box = slide.shapes.add_textbox(left_inch_text, top_inch_text, width_text, height_text)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.add_paragraph()
        p.text = create_slide_contents(slide_title)
        p.font.size = CONTENT_FONT_SIZE

        # Generate DALL-E image based on slide content
        dalle_image = generate_dalle_image(create_slide_contents(slide_title))

        if dalle_image:
            # Add an image on the right side
            left_inch_img = Inches(5)  # Adjust the left margin as needed
            top_inch_img = Inches(2.1)  # Adjust the top margin as needed
            width_img = Inches(4.5)  # Adjust the width of the image box
            height_img = Inches(4.5)  # Adjust the height of the image box
            img_stream = BytesIO()
            dalle_image.save(img_stream, format='PNG')
            pic = slide.shapes.add_picture(img_stream, left=left_inch_img, top=top_inch_img, width=width_img, height=height_img)

    powerpoint.save(f'powerpoint-ppt/{topic}.pptx')

    return f'powerpoint-ppt/{topic}.pptx'

######################################## FUNCTION CALLING WITH ASSISTANT API ###################################################################

assistant = client.beta.assistants.create(
  instructions="You are an expert in creating highly engaging powerpoint presentations",
  model="gpt-4-1106-preview",
  tools=[{
        "type": "function",
        "function": {
        "name": "create_presentations",
        "description": "Create presentation for the slide based on the topic & number of slides required passed as a parameter",
        "parameters": {
        "type": "object",
        "properties": {
                        "topic": {
                            "type": "string",
                            "description": "slide title, e.g. United Nations & It's Role",
                        },
                        "num_slides":{
                            "type":"string",
                            "description":"Number Of Slides Required"
                        }
                    },
            "required": ["topic","num_slides"]
        }
    }
  }]
)

thread = client.beta.threads.create()

message = client.beta.threads.messages.create(
    thread_id=thread.id,
    role="user",
    content="Can you please create presentations for the title Frontend Development & required number of slides is 5 ?"
)

run = client.beta.threads.runs.create(
  thread_id=thread.id,
  assistant_id=assistant.id,
  instructions=""
)

while True:
    # Wait for 5 seconds
    time.sleep(60)

    # Retrieve the run status
    run_status = client.beta.threads.runs.retrieve(
        thread_id=thread.id,
        run_id=run.id
    )
    # print(run_status.model_dump_json(indent=4))

    print("TIME = ",time.perf_counter())

    # If run is completed, get messages
    if run_status.status == 'completed':
        messages = client.beta.threads.messages.list(
            thread_id=thread.id
        )

        # Loop through messages and print content based on role
        for msg in messages.data:
            role = msg.role
            content = msg.content[0].text.value
            print(f"{role.capitalize()}: {content}")

        break
    elif run_status.status == 'requires_action':
        print("Function Calling")
        required_actions = run_status.required_action.submit_tool_outputs.model_dump()
        print("\nRequired Action\n\n",required_actions)
        tool_outputs = []
        for action in required_actions["tool_calls"]:
            func_name = action['function']['name']
            arguments = json.loads(action['function']['arguments'])
            
            if func_name == "create_presentations":
                print("\nArguments\n\n",arguments['topic'], arguments['num_slides'])
                output = create_presentations(arguments['topic'], arguments['num_slides'])
                tool_outputs.append({
                    "tool_call_id": action['id'],
                    "output": output
                })
                # time.sleep(60)
            else:
                raise ValueError(f"Unknown function: {func_name}")
            
        print("Submitting outputs back to the Assistant...")
        client.beta.threads.runs.submit_tool_outputs(
            thread_id=thread.id,
            run_id=run.id,
            tool_outputs=tool_outputs
        )
    else:
        print("Waiting for the Assistant to process...")
        time.sleep(60)
